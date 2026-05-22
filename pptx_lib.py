"""Helpers for building decks from template.pptx."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

# Layout indices in template.pptx (German names from slide master)
LAYOUT_TITLE = 0          # Titelfolie
LAYOUT_AGENDA = 2         # Agenda
LAYOUT_CHAPTER = 3        # Chapter Slide
LAYOUT_TITLE_CONTENT = 6  # Titel und Inhalt
LAYOUT_TWO_COLUMNS = 7    # Two Text Areas
LAYOUT_END = 33           # Logo Slide / End Slide

PROJECT_ROOT = Path(__file__).resolve().parent
DEFAULT_TEMPLATE = PROJECT_ROOT / "template.pptx"
DEFAULT_OUTPUT = PROJECT_ROOT / "output" / "Lernwoche-2026-Prompting-SW.pptx"


def load_template(path: Path | None = None) -> Presentation:
    template = path or DEFAULT_TEMPLATE
    if not template.is_file():
        raise FileNotFoundError(f"template.pptx not found: {template}")
    return Presentation(str(template))


def _placeholder_by_type(slide, ph_type: PP_PLACEHOLDER):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == ph_type:
            return shape
    return None


def _placeholder_by_idx(slide, idx: int):
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def set_center_title(slide, title: str, subtitle: str | None = None) -> None:
    title_ph = _placeholder_by_type(slide, PP_PLACEHOLDER.CENTER_TITLE) or _placeholder_by_idx(slide, 0)
    if title_ph is not None:
        title_ph.text = title
    if subtitle:
        sub_ph = _placeholder_by_type(slide, PP_PLACEHOLDER.SUBTITLE) or _placeholder_by_idx(slide, 1)
        if sub_ph is not None:
            sub_ph.text = subtitle


def set_slide_title(slide, title: str) -> None:
    title_ph = _placeholder_by_type(slide, PP_PLACEHOLDER.TITLE) or _placeholder_by_idx(slide, 0)
    if title_ph is not None:
        title_ph.text = title


def _clear_text_frame(tf) -> None:
    tf.clear()
    if not tf.paragraphs:
        tf.add_paragraph()


def set_bullets(shape, items: Iterable[str], font_size_pt: int = 18) -> None:
    tf = shape.text_frame
    _clear_text_frame(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size_pt)


def set_body_bullets(slide, items: Iterable[str], placeholder_idx: int = 13, font_size_pt: int = 18) -> None:
    body = _placeholder_by_idx(slide, placeholder_idx)
    if body is None:
        body = _placeholder_by_type(slide, PP_PLACEHOLDER.BODY)
    if body is None:
        raise RuntimeError("No body placeholder found on slide")
    set_bullets(body, items, font_size_pt)


def set_two_column_bullets(
    slide,
    left_items: Iterable[str],
    right_items: Iterable[str],
    left_idx: int = 13,
    right_idx: int = 14,
) -> None:
    left = _placeholder_by_idx(slide, left_idx)
    right = _placeholder_by_idx(slide, right_idx)
    if left is None or right is None:
        raise RuntimeError("Two-column placeholders not found (layout 7 expected)")
    set_bullets(left, left_items)
    set_bullets(right, right_items)


def add_slide(prs: Presentation, layout_index: int):
    return prs.slides.add_slide(prs.slide_layouts[layout_index])


def remove_template_sample_slides(prs: Presentation) -> None:
    """Remove all example slides shipped in template.pptx (count may change when template is updated)."""
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId  # noqa: SLF001
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def save_presentation(prs: Presentation, output_path: Path | None = None) -> Path:
    out = output_path or DEFAULT_OUTPUT
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
