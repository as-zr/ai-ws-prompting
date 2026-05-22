"""Custom slide builder — no corporate template."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0A, 0x2E, 0x4A)
NAVY_LIGHT = RGBColor(0x12, 0x45, 0x6E)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ORANGE_SOFT = RGBColor(0xFF, 0x9A, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF9, 0xFC)
INK = RGBColor(0x1A, 0x24, 0x32)
MUTED = RGBColor(0x5C, 0x6B, 0x7A)
BAD_BG = RGBColor(0xFF, 0xEB, 0xEE)
BAD_ACCENT = RGBColor(0xC6, 0x28, 0x28)
GOOD_BG = RGBColor(0xE8, 0xF5, 0xE9)
GOOD_ACCENT = RGBColor(0x2E, 0x7D, 0x32)
EXERCISE_BG = RGBColor(0xFF, 0xF8, 0xF0)

FONT = "Segoe UI"

PROJECT_ROOT = Path(__file__).resolve().parent
DEFAULT_OUTPUT = PROJECT_ROOT / "output" / "Lernwoche-2026-Prompting-SW.pptx"


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def _text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _bullets(
    slide,
    left,
    top,
    width,
    height,
    items: Iterable[str],
    *,
    size: int = 18,
    color: RGBColor = INK,
    spacing: int = 10,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return box


def slide_hero(
    prs: Presentation,
    title: str,
    subtitle: str,
    meta: str,
    tag: str = "Lernwoche 2026",
) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _rect(s, 0, Inches(6.85), SLIDE_W, Inches(0.12), ORANGE)
    _rect(s, Inches(0.85), Inches(2.0), Inches(0.08), Inches(2.8), ORANGE)

    _text(s, Inches(1.15), Inches(1.0), Inches(5), Inches(0.4), tag.upper(), size=11, bold=True, color=ORANGE_SOFT)
    _text(s, Inches(1.15), Inches(2.05), Inches(11), Inches(1.6), title, size=44, bold=True, color=WHITE)
    _text(s, Inches(1.15), Inches(3.55), Inches(10), Inches(0.9), subtitle, size=24, color=RGBColor(0xB8, 0xD4, 0xE8))
    _text(s, Inches(1.15), Inches(5.9), Inches(10), Inches(0.5), meta, size=15, color=MUTED)


def slide_agenda(prs: Presentation, title: str, items: list[tuple[str, str]]) -> None:
    """items: (label, time_or_hint)"""
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, OFF_WHITE)
    _rect(s, 0, 0, SLIDE_W, Inches(0.1), ORANGE)
    _text(s, Inches(0.9), Inches(0.55), Inches(8), Inches(0.8), title, size=34, bold=True, color=NAVY)

    y = Inches(1.55)
    for i, (label, hint) in enumerate(items, 1):
        _rect(s, Inches(0.9), y, Inches(0.45), Inches(0.45), ORANGE if i % 2 else NAVY_LIGHT)
        _text(
            s, Inches(0.9), y, Inches(0.45), Inches(0.45),
            str(i), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        _text(s, Inches(1.55), y - Inches(0.02), Inches(8.5), Inches(0.45), label, size=20, bold=True, color=INK)
        _text(s, Inches(10.2), y, Inches(2.5), Inches(0.45), hint, size=14, color=MUTED, align=PP_ALIGN.RIGHT)
        y += Inches(0.72)


def slide_chapter(prs: Presentation, chapter: str, headline: str) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_LIGHT)
    _rect(s, Inches(0.9), Inches(2.8), Inches(1.2), Inches(0.06), ORANGE)
    _text(s, Inches(0.9), Inches(1.6), Inches(11), Inches(0.6), chapter.upper(), size=14, bold=True, color=ORANGE_SOFT)
    _text(s, Inches(0.9), Inches(3.0), Inches(11), Inches(1.5), headline, size=40, bold=True, color=WHITE)


def slide_insight(prs: Presentation, title: str, lines: list[str], footer: str | None = None) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, OFF_WHITE)
    _rect(s, 0, 0, Inches(0.12), SLIDE_H, NAVY)
    _text(s, Inches(0.65), Inches(0.5), Inches(11), Inches(0.7), title, size=30, bold=True, color=NAVY)

    y = Inches(1.45)
    for line in lines:
        _rect(s, Inches(0.65), y + Inches(0.12), Inches(0.06), Inches(0.35), ORANGE)
        _text(s, Inches(0.9), y, Inches(11.5), Inches(0.7), line, size=22, color=INK)
        y += Inches(1.05)

    if footer:
        _text(s, Inches(0.65), Inches(6.5), Inches(11), Inches(0.5), footer, size=14, color=MUTED)


def slide_craft(prs: Presentation) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, OFF_WHITE)
    _text(s, Inches(0.75), Inches(0.45), Inches(11), Inches(0.6), "CRAFT — Prompt-Gerüst für technische Aufgaben", size=28, bold=True, color=NAVY)

    cards = [
        ("C", "Context / Kontext", "Dateien, Modul, Stack, aktueller Stand"),
        ("R", "Role / Persona", "Implementer, Reviewer, Tester, Teacher"),
        ("A", "Action / Aufgabe", "Konkreter Auftrag mit klarem Scope"),
        ("F", "Format / Ausgabe", "Diff, Schritte, Tests, Risiken, Sprache"),
        ("T", "Constraints / Grenzen", "Keine API-Änderung · keine Dependencies · Annahmen listen"),
    ]
    x0, y0 = Inches(0.75), Inches(1.35)
    w, h, gap = Inches(2.35), Inches(2.35), Inches(0.2)
    for i, (letter, name, desc) in enumerate(cards):
        col, row = i % 3, i // 3
        x = x0 + col * (w + gap)
        y = y0 + row * (h + gap)
        _rect(s, x, y, w, h, WHITE)
        _rect(s, x, y, w, Inches(0.08), ORANGE if i < 3 else NAVY)
        _text(s, x + Inches(0.15), y + Inches(0.2), Inches(0.6), Inches(0.7), letter, size=36, bold=True, color=ORANGE if i < 3 else NAVY)
        _text(s, x + Inches(0.15), y + Inches(0.85), w - Inches(0.3), Inches(0.4), name, size=16, bold=True, color=NAVY)
        _text(s, x + Inches(0.15), y + Inches(1.25), w - Inches(0.3), Inches(1.0), desc, size=13, color=MUTED)

    _text(
        s, Inches(0.75), Inches(6.35), Inches(11.8), Inches(0.5),
        "Anti-Patterns: zu vage · kein Kontext · keine Constraints · kein Review",
        size=14, bold=True, color=BAD_ACCENT,
    )


def slide_compare(prs: Presentation) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, OFF_WHITE)
    _text(s, Inches(0.75), Inches(0.45), Inches(11), Inches(0.6), "Gleicher Code — zwei Prompts", size=30, bold=True, color=NAVY)

    col_w = Inches(5.85)
    for left, bg, accent, label, prompt, bullets in [
        (Inches(0.75), BAD_BG, BAD_ACCENT, "SCHLECHT", "„Mach das schneller.“", ["Kein Kontext", "Keine Constraints", "Unklares Ziel"]),
        (
            Inches(6.75),
            GOOD_BG,
            GOOD_ACCENT,
            "GUT",
            "„Optimiere die DB-Abfrage in userRepository, Ziel: P95 < 150ms bei 10k Datensätzen, ohne API-Vertrag zu ändern.“",
            ["Präzises Ziel", "Messbare Kriterien", "Randbedingung klar"],
        ),
    ]:
        _rect(s, left, Inches(1.25), col_w, Inches(5.5), bg)
        _rect(s, left, Inches(1.25), col_w, Inches(0.55), accent)
        _text(s, left + Inches(0.2), Inches(1.32), col_w, Inches(0.4), label, size=14, bold=True, color=WHITE)
        _text(s, left + Inches(0.25), Inches(2.0), col_w - Inches(0.5), Inches(1.0), f"„{prompt}“", size=15, color=INK)
        _bullets(s, left + Inches(0.25), Inches(3.2), col_w - Inches(0.5), Inches(2.5), bullets, size=14, color=MUTED)

    _text(s, Inches(0.75), Inches(6.85), Inches(11.8), Inches(0.4), "Live: lesen → Annahmen prüfen → nachschärfen", size=13, color=NAVY, bold=True)


def slide_exercise(
    prs: Presentation,
    number: int,
    title: str,
    tasks: list[str],
    meta: str,
) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, EXERCISE_BG)
    _rect(s, 0, 0, Inches(0.35), SLIDE_H, ORANGE)
    _text(s, Inches(0.65), Inches(0.45), Inches(2), Inches(0.5), f"ÜBUNG {number}", size=13, bold=True, color=ORANGE)
    _text(s, Inches(0.65), Inches(0.85), Inches(11), Inches(0.8), title, size=32, bold=True, color=NAVY)
    _rect(s, Inches(0.65), Inches(1.75), Inches(2.2), Inches(0.4), ORANGE)
    _text(s, Inches(0.75), Inches(1.78), Inches(2.0), Inches(0.35), meta, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _bullets(s, Inches(0.65), Inches(2.45), Inches(11.5), Inches(4.0), tasks, size=19, color=INK, spacing=14)


def slide_workflow(prs: Presentation) -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, OFF_WHITE)
    _text(s, Inches(0.75), Inches(0.45), Inches(11), Inches(0.6), "Hands-on Workflow für reale Entwickler-Szenarien", size=28, bold=True, color=NAVY)

    steps = [
        ("1", "Markieren", "Methode im Editor wählen"),
        ("2", "Chat", "Copilot Chat öffnen"),
        ("3", "CRAFT", "Strukturierten Prompt senden"),
        ("4", "Review", "Iterieren & Tests laufen lassen"),
    ]
    x = Inches(0.55)
    step_w = Inches(2.95)
    for num, title, sub in steps:
        _rect(s, x, Inches(2.0), step_w, Inches(3.8), WHITE)
        _rect(s, x, Inches(2.0), step_w, Inches(0.06), ORANGE)
        _text(s, x + Inches(0.2), Inches(2.25), Inches(0.6), Inches(0.6), num, size=28, bold=True, color=ORANGE)
        _text(s, x + Inches(0.2), Inches(2.95), step_w - Inches(0.4), Inches(0.5), title, size=18, bold=True, color=NAVY)
        _text(s, x + Inches(0.2), Inches(3.5), step_w - Inches(0.4), Inches(1.5), sub, size=13, color=MUTED)
        if x < Inches(10):
            _text(s, x + step_w - Inches(0.15), Inches(3.6), Inches(0.4), Inches(0.4), "→", size=22, color=ORANGE)
        x += step_w + Inches(0.2)

    _text(
        s, Inches(0.75), Inches(6.2), Inches(11.8), Inches(0.8),
        "Use Cases: Implementierung · Refactoring · Tests · Debugging · Doku",
        size=13, color=MUTED,
    )


def slide_takeaway(prs: Presentation, items: list[tuple[str, str]]) -> None:
    """(title, subtitle) per column."""
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _text(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.7), "Mitnehmen für Montag", size=34, bold=True, color=WHITE)

    n = len(items)
    col_w = Inches(11.5 / n) if n else Inches(3.8)
    x = Inches(0.9)
    for title, sub in items:
        _rect(s, x, Inches(2.0), col_w - Inches(0.25), Inches(4.2), NAVY_LIGHT)
        _rect(s, x, Inches(2.0), col_w - Inches(0.25), Inches(0.06), ORANGE)
        _text(s, x + Inches(0.2), Inches(2.35), col_w - Inches(0.65), Inches(1.2), title, size=20, bold=True, color=WHITE)
        _text(s, x + Inches(0.2), Inches(3.6), col_w - Inches(0.65), Inches(2.0), sub, size=14, color=RGBColor(0xB8, 0xD4, 0xE8))
        x += col_w


def slide_end(prs: Presentation, line: str = "Fragen?") -> None:
    s = _blank(prs)
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _rect(s, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.08), ORANGE)
    _text(s, Inches(0.9), Inches(3.45), Inches(11.5), Inches(1.0), line, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _text(
        s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.5),
        "Lernwoche 2026 · GitHub Copilot · 90 Min Hands-on",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )


def save(prs: Presentation, path: Path | None = None) -> Path:
    out = path or DEFAULT_OUTPUT
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
